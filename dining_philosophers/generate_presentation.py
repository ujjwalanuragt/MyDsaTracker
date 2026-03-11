"""
Dining Philosophers Problem — PowerPoint Presentation Generator
Generates a 13-slide academic PPTX using python-pptx (no external images).
Run: python generate_presentation.py
Output: Dining_Philosophers_Problem.pptx
"""

import math
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour palette ──────────────────────────────────────────────────────────
BG_COLOR       = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
GOLD           = RGBColor(0xF5, 0xA6, 0x23)   # Golden yellow
STEEL_BLUE     = RGBColor(0xB0, 0xC4, 0xDE)   # Light steel blue
DARK_BLUE      = RGBColor(0x1A, 0x35, 0x50)   # Card background
CODE_BG        = RGBColor(0x0A, 0x15, 0x20)   # Code block background
RED            = RGBColor(0xFF, 0x44, 0x44)    # Deadlock arrows

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide):
    """Set slide background to deep navy."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=18, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a plain textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_rect(slide, left, top, width, height,
             fill_color=DARK_BLUE, line_color=None, line_width=Pt(1)):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=DARK_BLUE, line_color=GOLD, line_width=Pt(1.5)):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        5,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_circle(slide, cx, cy, diameter,
               fill_color=GOLD, line_color=None):
    """Add a filled circle centred at (cx, cy)."""
    left = cx - diameter / 2
    top  = cy - diameter / 2
    shape = slide.shapes.add_shape(
        9,  # MSO_AUTO_SHAPE_TYPE.OVAL
        int(left), int(top), int(diameter), int(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_slide_title(slide, title_text, top=Inches(0.25)):
    """Add a golden slide title."""
    add_textbox(slide, title_text,
                left=Inches(0.4), top=top,
                width=Inches(12.5), height=Inches(0.7),
                font_size=30, bold=True, color=GOLD,
                align=PP_ALIGN.LEFT)


def philosopher_positions(cx, cy, radius, n=5):
    """Return (x, y) centres for n items in a circle, starting at top."""
    positions = []
    for i in range(n):
        angle = math.radians(-90 + i * 360 / n)
        x = cx + radius * math.cos(angle)
        y = cy + radius * math.sin(angle)
        positions.append((x, y))
    return positions


def chopstick_positions(cx, cy, radius, n=5):
    """Return midpoints between adjacent philosopher positions."""
    phil = philosopher_positions(cx, cy, radius, n)
    mids = []
    for i in range(n):
        x = (phil[i][0] + phil[(i + 1) % n][0]) / 2
        y = (phil[i][1] + phil[(i + 1) % n][1]) / 2
        mids.append((x, y))
    return mids


def draw_table_diagram(slide, cx, cy, table_r,
                       phil_r, phil_dot=Inches(0.28),
                       chop_size=(Inches(0.18), Inches(0.08)),
                       show_deadlock=False):
    """
    Draw the circular-table dining-philosophers diagram.
    cx, cy   – centre of the diagram (Emu)
    table_r  – radius of the table circle (Emu)
    phil_r   – orbit radius for philosopher circles (Emu)
    """
    # Table
    add_circle(slide, cx, cy, table_r * 2,
               fill_color=DARK_BLUE, line_color=GOLD)
    add_textbox(slide, "Table",
                left=int(cx - Inches(0.5)), top=int(cy - Inches(0.2)),
                width=Inches(1), height=Inches(0.4),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Philosopher circles + labels
    phil_pos = philosopher_positions(cx, cy, phil_r)
    for i, (px, py) in enumerate(phil_pos):
        fill = GOLD if not show_deadlock else RGBColor(0xFF, 0xA0, 0x00)
        add_circle(slide, px, py, phil_dot, fill_color=fill)
        # label
        lx = px + (px - cx) * 0.32
        ly = py + (py - cy) * 0.32
        add_textbox(slide, f"P{i}",
                    left=int(lx - Inches(0.22)), top=int(ly - Inches(0.18)),
                    width=Inches(0.44), height=Inches(0.36),
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Chopstick rectangles
    chop_pos = chopstick_positions(cx, cy, (table_r + phil_r) / 2)
    chop_w, chop_h = chop_size
    for i, (cpx, cpy) in enumerate(chop_pos):
        fill = GOLD if show_deadlock else WHITE
        add_rect(slide,
                 left=int(cpx - chop_w / 2), top=int(cpy - chop_h / 2),
                 width=int(chop_w), height=int(chop_h),
                 fill_color=fill, line_color=None)
        add_textbox(slide, f"C{i}",
                    left=int(cpx - Inches(0.22)),
                    top=int(cpy - Inches(0.32)),
                    width=Inches(0.44), height=Inches(0.26),
                    font_size=9, color=STEEL_BLUE, align=PP_ALIGN.CENTER)

    if show_deadlock:
        # Dashed red arrows indicating "waiting for" (visual only — use thin lines)
        for i, (px, py) in enumerate(phil_pos):
            # next chopstick (clockwise)
            next_chop = chop_pos[i]
            _draw_arrow_line(slide, px, py, next_chop[0], next_chop[1], RED)


def _draw_arrow_line(slide, x1, y1, x2, y2, color):
    """Draw a simple line connector (visual arrow substitute)."""
    # Use a connector shape as a thin line
    connector = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_slide1_title(prs):
    """Slide 1 – Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    W = prs.slide_width
    H = prs.slide_height

    # Title
    add_textbox(slide, "Dining Philosophers Problem",
                left=Inches(0.6), top=Inches(1.4),
                width=Inches(12.1), height=Inches(1.2),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, "Process Synchronization in Operating Systems",
                left=Inches(0.6), top=Inches(2.7),
                width=Inches(12.1), height=Inches(0.7),
                font_size=22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    # Decorative golden horizontal rule
    rule = slide.shapes.add_shape(1,
        Inches(1.5), Inches(3.55), Inches(10.33), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = GOLD
    rule.line.fill.background()

    # Info block
    info_lines = [
        "Presented by: Ujjwal Anuragt",
        "Course: Operating Systems",
        "Date: March 11, 2026",
    ]
    for idx, line in enumerate(info_lines):
        add_textbox(slide, line,
                    left=Inches(1.5), top=Inches(3.75 + idx * 0.48),
                    width=Inches(9), height=Inches(0.44),
                    font_size=18, color=STEEL_BLUE, align=PP_ALIGN.CENTER)

    # 5 small circles (pentagon) — bottom-right corner hint
    corner_cx = Inches(11.8)
    corner_cy = Inches(6.5)
    corner_r  = Inches(0.55)
    for i in range(5):
        angle = math.radians(-90 + i * 72)
        px = corner_cx + corner_r * math.cos(angle)
        py = corner_cy + corner_r * math.sin(angle)
        add_circle(slide, int(px), int(py), int(Inches(0.16)), fill_color=GOLD)


def build_slide2_introduction(prs):
    """Slide 2 – Introduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_slide_title(slide, "Introduction")

    # Left block
    add_rounded_rect(slide, Inches(0.4), Inches(1.1), Inches(5.5), Inches(2.2),
                     fill_color=DARK_BLUE, line_color=GOLD)
    add_textbox(slide, "Process Synchronization",
                left=Inches(0.55), top=Inches(1.15),
                width=Inches(5.2), height=Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide,
                "Ensures orderly access to shared resources by concurrent "
                "processes to prevent conflicts.",
                left=Inches(0.55), top=Inches(1.65),
                width=Inches(5.2), height=Inches(1.5),
                font_size=14, color=STEEL_BLUE)

    # Right block
    add_rounded_rect(slide, Inches(6.2), Inches(1.1), Inches(5.5), Inches(2.2),
                     fill_color=DARK_BLUE, line_color=GOLD)
    add_textbox(slide, "Resource Sharing",
                left=Inches(6.35), top=Inches(1.15),
                width=Inches(5.2), height=Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide,
                "Multiple processes compete for limited resources, requiring "
                "careful coordination.",
                left=Inches(6.35), top=Inches(1.65),
                width=Inches(5.2), height=Inches(1.5),
                font_size=14, color=STEEL_BLUE)

    # Bottom full-width block
    add_rounded_rect(slide, Inches(0.4), Inches(3.55), Inches(9.5), Inches(1.5),
                     fill_color=DARK_BLUE, line_color=STEEL_BLUE)
    add_textbox(slide,
                "The Dining Philosophers Problem was introduced by Edsger Dijkstra "
                "in 1965 to model synchronization issues such as deadlock and "
                "starvation in concurrent systems.",
                left=Inches(0.55), top=Inches(3.65),
                width=Inches(9.2), height=Inches(1.3),
                font_size=14, color=WHITE)

    # Person icon (right margin)
    # Body rectangle
    add_rect(slide, Inches(10.5), Inches(3.7), Inches(1.4), Inches(1.8),
             fill_color=DARK_BLUE, line_color=STEEL_BLUE)
    # Head circle
    add_circle(slide, int(Inches(11.2)), int(Inches(3.55)), int(Inches(0.4)),
               fill_color=STEEL_BLUE)
    add_textbox(slide, "Edsger\nDijkstra",
                left=Inches(10.5), top=Inches(3.85),
                width=Inches(1.4), height=Inches(0.7),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "1965",
                left=Inches(10.5), top=Inches(5.0),
                width=Inches(1.4), height=Inches(0.35),
                font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def build_slide3_overview(prs):
    """Slide 3 – Problem Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Problem Overview")

    bullets = [
        "A classical thought experiment in computer science.",
        "Models the challenge of allocating shared resources among competing processes.",
        "Highlights:  Deadlock  ·  Starvation  ·  Mutual Exclusion",
        "Originally stated with 5 philosophers, 5 forks/chopsticks, and a circular table.",
        "Widely used as a benchmark to evaluate synchronization algorithms.",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide, f"• {b}",
                    left=Inches(0.4), top=Inches(1.1 + i * 0.58),
                    width=Inches(7.5), height=Inches(0.52),
                    font_size=16, color=STEEL_BLUE)

    # State cycle diagram (Think → Hungry → Eat → Think)
    states = ["Think", "Hungry", "Eat"]
    state_x = [Inches(9.2), Inches(11.0), Inches(10.1)]
    state_y = [Inches(1.6),  Inches(2.9),  Inches(4.5)]

    for sx, sy, label in zip(state_x, state_y, states):
        add_rounded_rect(slide, int(sx - Inches(0.55)), int(sy - Inches(0.25)),
                         Inches(1.1), Inches(0.5),
                         fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, label,
                    left=int(sx - Inches(0.55)), top=int(sy - Inches(0.25)),
                    width=Inches(1.1), height=Inches(0.5),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows between states
    pairs = [(0, 1), (1, 2), (2, 0)]
    for a, b_idx in pairs:
        x1 = state_x[a]
        y1 = state_y[a] + Inches(0.25)
        x2 = state_x[b_idx]
        y2 = state_y[b_idx] - Inches(0.25)
        _draw_arrow_line(slide, x1, y1, x2, y2, GOLD)


def build_slide4_scenario(prs):
    """Slide 4 – Problem Scenario."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Problem Scenario")

    bullets = [
        "5 philosophers seated at a circular table",
        "1 chopstick placed between each adjacent pair",
        "Philosophers alternate between thinking and eating",
        "Eating requires BOTH left and right chopsticks",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide, f"• {b}",
                    left=Inches(0.4), top=Inches(1.1 + i * 0.65),
                    width=Inches(6.5), height=Inches(0.58),
                    font_size=16, color=STEEL_BLUE)

    # Circular table diagram (right half)
    draw_table_diagram(slide,
                       cx=Inches(10.0), cy=Inches(4.2),
                       table_r=Inches(1.1), phil_r=Inches(2.0))


def build_slide5_resource_sharing(prs):
    """Slide 5 – Resource Sharing Representation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Resource Sharing Representation")

    content = [
        "Philosophers  →  Processes (threads competing for CPU/resources)",
        "Chopsticks  →  Shared Resources (files, memory, devices)",
        "Thinking  →  Process is in non-critical section",
        "Eating  →  Process is in critical section using shared resources",
    ]
    for i, line in enumerate(content):
        add_textbox(slide, f"• {line}",
                    left=Inches(0.4), top=Inches(1.1 + i * 0.7),
                    width=Inches(6.8), height=Inches(0.6),
                    font_size=15, color=STEEL_BLUE)

    # Mapping diagram (right side)
    procs = ["Process 1", "Process 2", "Process 3"]
    res   = ["Resource A", "Resource B"]
    px_left  = Inches(8.2)
    px_right = Inches(10.8)
    box_w    = Inches(1.8)
    box_h    = Inches(0.5)

    proc_y = [Inches(1.5 + i * 1.1) for i in range(3)]
    res_y  = [Inches(2.1 + i * 1.5) for i in range(2)]

    for i, (label, py) in enumerate(zip(procs, proc_y)):
        add_rect(slide, int(px_left), int(py), int(box_w), int(box_h),
                 fill_color=DARK_BLUE, line_color=STEEL_BLUE)
        add_textbox(slide, label,
                    left=int(px_left), top=int(py),
                    width=int(box_w), height=int(box_h),
                    font_size=13, color=STEEL_BLUE, align=PP_ALIGN.CENTER)

    for i, (label, ry) in enumerate(zip(res, res_y)):
        add_rect(slide, int(px_right), int(ry), int(box_w), int(box_h),
                 fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, label,
                    left=int(px_right), top=int(ry),
                    width=int(box_w), height=int(box_h),
                    font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

    # Arrows: proc → resource
    arrow_pairs = [(0, 0, "holds"), (1, 0, "requests"), (1, 1, "holds"), (2, 1, "requests")]
    for pi, ri, lbl in arrow_pairs:
        x1 = int(px_left  + box_w)
        y1 = int(proc_y[pi] + box_h / 2)
        x2 = int(px_right)
        y2 = int(res_y[ri] + box_h / 2)
        _draw_arrow_line(slide, x1, y1, x2, y2, STEEL_BLUE)
        mid_x = (x1 + x2) // 2
        mid_y = (y1 + y2) // 2
        add_textbox(slide, lbl,
                    left=mid_x - int(Inches(0.35)), top=mid_y - int(Inches(0.2)),
                    width=int(Inches(0.7)), height=int(Inches(0.28)),
                    font_size=9, color=GOLD, align=PP_ALIGN.CENTER)


def build_slide6_challenges(prs):
    """Slide 6 – Synchronization Challenges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Synchronization Challenges")

    cards = [
        ("Race Conditions",
         "Multiple processes accessing shared resources simultaneously "
         "without coordination leads to unpredictable results."),
        ("Ordering Constraints",
         "Certain operations must execute in a specific order across "
         "concurrent processes."),
        ("Resource Contention",
         "Processes compete for limited shared resources, causing delays "
         "or conflicts."),
    ]

    card_w  = Inches(3.9)
    card_h  = Inches(4.8)
    gap     = Inches(0.25)
    start_x = Inches(0.4)
    card_top = Inches(1.05)

    for i, (heading, body) in enumerate(cards):
        cx = start_x + i * (card_w + gap)
        # Card background
        add_rounded_rect(slide, int(cx), int(card_top), int(card_w), int(card_h),
                         fill_color=DARK_BLUE, line_color=DARK_BLUE)
        # Gold top border accent
        accent = slide.shapes.add_shape(1,
            int(cx), int(card_top), int(card_w), int(Inches(0.12)))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD
        accent.line.fill.background()
        # Heading
        add_textbox(slide, heading,
                    left=int(cx + Inches(0.15)), top=int(card_top + Inches(0.2)),
                    width=int(card_w - Inches(0.3)), height=Inches(0.5),
                    font_size=17, bold=True, color=WHITE)
        # Body
        add_textbox(slide, body,
                    left=int(cx + Inches(0.15)), top=int(card_top + Inches(0.82)),
                    width=int(card_w - Inches(0.3)), height=Inches(3.7),
                    font_size=14, color=STEEL_BLUE)


def build_slide7_key_issues(prs):
    """Slide 7 – Key Issues."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Key Issues")

    sections = [
        ("Mutual Exclusion", [
            "Only one philosopher can use a chopstick at a time.",
            "Enforced to prevent simultaneous resource access.",
        ]),
        ("Deadlock", [
            "All philosophers pick up the left chopstick and wait for the right.",
            "Circular wait — no progress is made.",
        ]),
        ("Starvation", [
            "A philosopher waits indefinitely while others repeatedly eat.",
            "Fairness is not guaranteed even without deadlock.",
        ]),
    ]

    sec_w   = Inches(5.5)
    sec_h   = Inches(2.0)
    sec_gap = Inches(0.18)
    sx_left = Inches(0.4)

    for i, (heading, points) in enumerate(sections):
        sy = Inches(1.1 + i * (sec_h + sec_gap))
        add_rounded_rect(slide, int(sx_left), int(sy), int(sec_w), int(sec_h),
                         fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, heading,
                    left=int(sx_left + Inches(0.15)), top=int(sy + Inches(0.08)),
                    width=int(sec_w - Inches(0.3)), height=Inches(0.42),
                    font_size=17, bold=True, color=GOLD)
        for j, pt in enumerate(points):
            add_textbox(slide, f"• {pt}",
                        left=int(sx_left + Inches(0.15)),
                        top=int(sy + Inches(0.55) + j * Inches(0.58)),
                        width=int(sec_w - Inches(0.3)), height=Inches(0.52),
                        font_size=14, color=STEEL_BLUE)

    # Deadlock circular-wait arrows (right side, row 2)
    dead_cx = Inches(10.3)
    dead_cy = Inches(4.3)
    dead_r  = Inches(1.4)
    n = 5
    pts = philosopher_positions(dead_cx, dead_cy, dead_r, n)
    for i in range(n):
        x1, y1 = pts[i]
        x2, y2 = pts[(i + 1) % n]
        _draw_arrow_line(slide, int(x1), int(y1), int(x2), int(y2), RED)
    add_textbox(slide, "Circular\nWait",
                left=int(dead_cx - Inches(0.6)), top=int(dead_cy - Inches(0.35)),
                width=Inches(1.2), height=Inches(0.7),
                font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)


def build_slide8_semaphore_solution(prs):
    """Slide 8 – Semaphore-Based Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Semaphore-Based Solution")

    # Definition box
    add_rounded_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.85),
                     fill_color=DARK_BLUE, line_color=GOLD, line_width=Pt(2))
    add_textbox(slide,
                "A semaphore is a synchronization variable used to control access "
                "to shared resources by multiple processes.",
                left=Inches(0.55), top=Inches(1.1),
                width=Inches(12.2), height=Inches(0.75),
                font_size=15, color=WHITE)

    # Content bullets
    bullets = [
        "Each chopstick is represented as a semaphore.",
        "Semaphore value = 1   →   chopstick is available.",
        "Semaphore value = 0   →   chopstick is in use.",
        "Declaration:   semaphore chopstick[5] = {1, 1, 1, 1, 1}",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide, f"• {b}",
                    left=Inches(0.4), top=Inches(2.1 + i * 0.55),
                    width=Inches(12.5), height=Inches(0.5),
                    font_size=16, color=STEEL_BLUE)

    # 5 semaphore circles at bottom
    add_textbox(slide, "Chopstick Semaphores",
                left=Inches(2.8), top=Inches(4.85),
                width=Inches(7.5), height=Inches(0.4),
                font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    for i in range(5):
        cx = Inches(3.8 + i * 1.3)
        cy = Inches(5.8)
        add_circle(slide, int(cx), int(cy), int(Inches(0.5)),
                   fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, "1",
                    left=int(cx - Inches(0.25)), top=int(cy - Inches(0.18)),
                    width=Inches(0.5), height=Inches(0.36),
                    font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"C{i}",
                    left=int(cx - Inches(0.25)), top=int(cy + Inches(0.3)),
                    width=Inches(0.5), height=Inches(0.3),
                    font_size=12, color=STEEL_BLUE, align=PP_ALIGN.CENTER)


def build_slide9_semaphore_ops(prs):
    """Slide 9 – Semaphore Operations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Semaphore Operations")

    # Left code box — wait()
    wait_code = (
        "wait(S):\n"
        "  while (S <= 0)\n"
        "    ; // busy wait\n"
        "  S--;"
    )
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(5.8), Inches(2.2),
             fill_color=CODE_BG, line_color=GOLD)
    add_textbox(slide, "wait()  operation",
                left=Inches(0.55), top=Inches(1.1),
                width=Inches(5.5), height=Inches(0.42),
                font_size=15, bold=True, color=GOLD)
    add_textbox(slide, wait_code,
                left=Inches(0.55), top=Inches(1.58),
                width=Inches(5.5), height=Inches(1.55),
                font_name="Courier New", font_size=14, color=WHITE)
    add_textbox(slide,
                "Decrements semaphore. Blocks process if S = 0.\nUsed to ACQUIRE a resource.",
                left=Inches(0.55), top=Inches(3.35),
                width=Inches(5.5), height=Inches(0.8),
                font_size=13, color=STEEL_BLUE)

    # Right code box — signal()
    sig_code = "signal(S):\n  S++;"
    add_rect(slide, Inches(7.0), Inches(1.05), Inches(5.8), Inches(2.2),
             fill_color=CODE_BG, line_color=GOLD)
    add_textbox(slide, "signal()  operation",
                left=Inches(7.15), top=Inches(1.1),
                width=Inches(5.5), height=Inches(0.42),
                font_size=15, bold=True, color=GOLD)
    add_textbox(slide, sig_code,
                left=Inches(7.15), top=Inches(1.58),
                width=Inches(5.5), height=Inches(1.55),
                font_name="Courier New", font_size=14, color=WHITE)
    add_textbox(slide,
                "Increments semaphore. Wakes a waiting process.\nUsed to RELEASE a resource.",
                left=Inches(7.15), top=Inches(3.35),
                width=Inches(5.5), height=Inches(0.8),
                font_size=13, color=STEEL_BLUE)

    # Flow diagram below
    flow_items = [
        (Inches(1.0),  Inches(4.6), "Process wants\nresource"),
        (Inches(3.8),  Inches(4.6), "S > 0 ?"),
        (Inches(6.6),  Inches(4.6), "Decrement S"),
        (Inches(9.0),  Inches(4.6), "Use Resource"),
        (Inches(11.3), Inches(4.6), "Increment S"),
        (Inches(3.8),  Inches(6.0), "Block / Wait"),
    ]
    box_w = Inches(1.6)
    box_h = Inches(0.55)

    for idx, (fx, fy, label) in enumerate(flow_items):
        color = GOLD if idx == 1 else DARK_BLUE
        lc    = GOLD if idx == 1 else STEEL_BLUE
        add_rounded_rect(slide, int(fx), int(fy), int(box_w), int(box_h),
                         fill_color=color, line_color=lc)
        tc = DARK_BLUE if idx == 1 else WHITE
        add_textbox(slide, label,
                    left=int(fx), top=int(fy),
                    width=int(box_w), height=int(box_h),
                    font_size=11, bold=(idx == 1), color=tc,
                    align=PP_ALIGN.CENTER)

    # Horizontal arrows
    for i in range(4):
        x1 = int(flow_items[i][0] + box_w)
        y1 = int(flow_items[i][1] + box_h / 2)
        x2 = int(flow_items[i + 1][0])
        y2 = int(flow_items[i + 1][1] + box_h / 2)
        _draw_arrow_line(slide, x1, y1, x2, y2, GOLD)

    # "No" branch down then arrow back to diamond
    no_x  = int(flow_items[1][0] + box_w / 2)
    no_y1 = int(flow_items[1][1] + box_h)
    no_y2 = int(flow_items[5][1])
    _draw_arrow_line(slide, no_x, no_y1, no_x, no_y2, RED)
    _draw_arrow_line(slide, int(flow_items[5][0] + box_w), int(flow_items[5][1] + box_h / 2),
                     no_x, int(flow_items[5][1] + box_h / 2), RED)


def build_slide10_pseudocode(prs):
    """Slide 10 – Pseudocode Using Semaphores."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Pseudocode — Semaphore Solution")

    code = (
        "semaphore chopstick[5] = {1, 1, 1, 1, 1};\n\n"
        "Philosopher(i):\n"
        "  while (true) {\n"
        "      think();\n\n"
        "      wait(chopstick[i]);             // Pick up LEFT chopstick\n"
        "      wait(chopstick[(i+1) % 5]);     // Pick up RIGHT chopstick\n\n"
        "      eat();\n\n"
        "      signal(chopstick[i]);           // Put down LEFT chopstick\n"
        "      signal(chopstick[(i+1) % 5]);   // Put down RIGHT chopstick\n"
        "  }"
    )

    # Code background
    add_rect(slide, Inches(0.35), Inches(1.0), Inches(9.2), Inches(6.2),
             fill_color=CODE_BG, line_color=GOLD)
    add_textbox(slide, code,
                left=Inches(0.5), top=Inches(1.1),
                width=Inches(8.9), height=Inches(6.0),
                font_name="Courier New", font_size=14, color=WHITE)

    # Annotation callouts (right column)
    annotations = [
        (Inches(1.55), "Acquire left\nchopstick"),
        (Inches(2.35), "Acquire right\nchopstick"),
        (Inches(3.35), "Critical\nSection"),
        (Inches(4.45), "Release both\nchopsticks"),
    ]
    ann_x = Inches(9.85)
    for ann_y, label in annotations:
        add_rounded_rect(slide, int(ann_x), int(ann_y - Inches(0.08)),
                         Inches(3.1), Inches(0.58),
                         fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, label,
                    left=int(ann_x + Inches(0.08)), top=int(ann_y - Inches(0.08)),
                    width=Inches(2.9), height=Inches(0.58),
                    font_size=12, color=GOLD, align=PP_ALIGN.CENTER)
        # Pointer line
        _draw_arrow_line(slide,
                         int(ann_x), int(ann_y + Inches(0.21)),
                         int(Inches(9.55)), int(ann_y + Inches(0.21)),
                         GOLD)


def build_slide11_deadlock_problem(prs):
    """Slide 11 – Problem with Basic Semaphore Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Problem: Deadlock in Basic Solution")

    bullets = [
        "If ALL 5 philosophers simultaneously pick up their LEFT chopstick:",
        "→ Each holds 1 chopstick and waits for the other.",
        "→ No philosopher can proceed to eat.",
        "→ DEADLOCK — circular wait condition.",
        "Coffman Conditions met: Mutual Exclusion, Hold & Wait,\n"
        "   No Preemption, Circular Wait.",
    ]
    for i, b in enumerate(bullets):
        col = GOLD if b.startswith("→") or b.startswith("Coffman") else STEEL_BLUE
        add_textbox(slide, b,
                    left=Inches(0.4), top=Inches(1.1 + i * 0.72),
                    width=Inches(7.0), height=Inches(0.65),
                    font_size=15, color=col)

    # Deadlock diagram (right side, showing all chopsticks golden = "held")
    draw_table_diagram(slide,
                       cx=Inches(10.5), cy=Inches(4.3),
                       table_r=Inches(1.0), phil_r=Inches(1.85),
                       show_deadlock=True)
    add_textbox(slide, "All chopsticks held!\nCircular wait →\nDEADLOCK",
                left=Inches(8.5), top=Inches(6.2),
                width=Inches(4.0), height=Inches(0.9),
                font_size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)


def build_slide12_prevention(prs):
    """Slide 12 – Deadlock Prevention Methods."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Deadlock Prevention Methods")

    solutions = [
        ("Solution 1 — Limit Philosophers",
         "Allow at most 4 philosophers to sit at the table simultaneously.\n"
         "At least one philosopher will always be able to eat."),
        ("Solution 2 — Atomic Pickup",
         "A philosopher picks up chopsticks ONLY if BOTH are available.\n"
         "Prevents the hold-and-wait condition."),
        ("Solution 3 — Asymmetric Solution",
         "Odd-numbered philosophers: Pick LEFT chopstick first, then RIGHT.\n"
         "Even-numbered philosophers: Pick RIGHT chopstick first, then LEFT.\n"
         "Breaks the symmetry that causes circular wait."),
    ]

    band_w  = Inches(7.5)
    band_h  = Inches(1.6)
    band_x  = Inches(0.4)

    for i, (heading, body) in enumerate(solutions):
        by = Inches(1.05 + i * (band_h + Inches(0.12)))
        add_rounded_rect(slide, int(band_x), int(by), int(band_w), int(band_h),
                         fill_color=DARK_BLUE, line_color=GOLD)
        # Number circle
        add_circle(slide,
                   int(band_x + Inches(0.38)), int(by + band_h / 2),
                   int(Inches(0.5)),
                   fill_color=GOLD)
        add_textbox(slide, str(i + 1),
                    left=int(band_x + Inches(0.13)), top=int(by + band_h / 2 - Inches(0.22)),
                    width=Inches(0.5), height=Inches(0.44),
                    font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, heading,
                    left=int(band_x + Inches(0.8)), top=int(by + Inches(0.08)),
                    width=int(band_w - Inches(0.95)), height=Inches(0.45),
                    font_size=15, bold=True, color=WHITE)
        add_textbox(slide, body,
                    left=int(band_x + Inches(0.8)), top=int(by + Inches(0.58)),
                    width=int(band_w - Inches(0.95)), height=Inches(0.95),
                    font_size=13, color=STEEL_BLUE)

    # Mini asymmetric diagram (right side)
    mini_cx = Inches(10.5)
    mini_cy = Inches(3.8)
    mini_r  = Inches(1.5)

    draw_table_diagram(slide,
                       cx=mini_cx, cy=mini_cy,
                       table_r=Inches(0.85), phil_r=mini_r,
                       phil_dot=Inches(0.22))

    add_textbox(slide, "P0 (even): RIGHT first\nP1 (odd):  LEFT first",
                left=Inches(8.8), top=Inches(5.6),
                width=Inches(4.0), height=Inches(0.7),
                font_size=12, color=STEEL_BLUE, align=PP_ALIGN.CENTER)


def build_slide13_conclusion(prs):
    """Slide 13 – Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_title(slide, "Conclusion")

    cards = [
        "Models synchronization and resource allocation issues in operating systems.",
        "Demonstrates deadlock and starvation in concurrent systems.",
        "Semaphore-based solutions effectively manage shared resource access.",
        "Additional strategies (limiting philosophers, asymmetric rules) are\n"
        "required to fully prevent deadlock.",
    ]

    card_w = Inches(5.9)
    card_h = Inches(2.2)
    positions = [
        (Inches(0.35), Inches(1.05)),
        (Inches(6.55), Inches(1.05)),
        (Inches(0.35), Inches(3.45)),
        (Inches(6.55), Inches(3.45)),
    ]

    for (cx, cy), text in zip(positions, cards):
        add_rounded_rect(slide, int(cx), int(cy), int(card_w), int(card_h),
                         fill_color=DARK_BLUE, line_color=GOLD)
        add_textbox(slide, text,
                    left=int(cx + Inches(0.15)), top=int(cy + Inches(0.2)),
                    width=int(card_w - Inches(0.3)), height=int(card_h - Inches(0.3)),
                    font_size=14, color=STEEL_BLUE)

    # Bottom banner
    banner = slide.shapes.add_shape(1,
        Inches(0.0), Inches(6.2), Inches(13.33), Inches(1.05))
    banner.fill.solid()
    banner.fill.fore_color.rgb = GOLD
    banner.line.fill.background()
    add_textbox(slide,
                "The Dining Philosophers Problem remains a foundational concept "
                "in Operating Systems and Concurrent Programming.",
                left=Inches(0.3), top=Inches(6.28),
                width=Inches(12.7), height=Inches(0.85),
                font_size=15, bold=True, color=BG_COLOR, align=PP_ALIGN.CENTER)

    # Small table diagram top-right
    draw_table_diagram(slide,
                       cx=Inches(12.0), cy=Inches(0.8),
                       table_r=Inches(0.45), phil_r=Inches(0.78),
                       phil_dot=Inches(0.14),
                       chop_size=(Inches(0.1), Inches(0.05)))


# ─────────────────────────────────────────────────────────────────────────────
#  MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    build_slide1_title(prs)
    build_slide2_introduction(prs)
    build_slide3_overview(prs)
    build_slide4_scenario(prs)
    build_slide5_resource_sharing(prs)
    build_slide6_challenges(prs)
    build_slide7_key_issues(prs)
    build_slide8_semaphore_solution(prs)
    build_slide9_semaphore_ops(prs)
    build_slide10_pseudocode(prs)
    build_slide11_deadlock_problem(prs)
    build_slide12_prevention(prs)
    build_slide13_conclusion(prs)

    out_path = os.path.join(os.path.dirname(__file__), "Dining_Philosophers_Problem.pptx")
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
